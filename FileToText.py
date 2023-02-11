from PyPDF2 import PdfReader
from docx2python import docx2python
from pptx import Presentation



class FileToText:

    # max string size in chars
    max_output_length = 1000

    @staticmethod
    def pdf_to_text(path, CUT_STR=False, max_output_length=max_output_length, GET_PAGES=True, pages_to_get=1,):
        DEBUG = True
        file_text = ""
        # note: only works with pdfs where text is not an image
        reader = PdfReader(f"{path}")

        # GET_PAGES bool specifies whether to extract text from the whole document
        if GET_PAGES:
            # in a given pdf doc
            number_of_pages = len(reader.pages)
        else:
            number_of_pages = pages_to_get

        for page in reader.pages:
            text = page.extract_text().replace("\n", " ")
            file_text += text

        #removes multiple whitespaces
        file_text = " ".join(file_text.split())

        if CUT_STR:
            file_text = file_text[:max_output_length]

        return file_text

    
    @staticmethod
    def docx_to_text(path, CUT_STR=False, max_output_length=max_output_length):
        file_text = ""

        with docx2python(f"{path}") as docx_content:
            file_text = docx_content.text.replace("\n", " ")
            
            #removes multiple whitespaces
            file_text = " ".join(file_text.split())

        if CUT_STR:
            file_text = file_text[:max_output_length]

        return file_text


    @staticmethod
    def pptx_to_text(path, CUT_STR=False, max_output_length=max_output_length):
        file_text = ""
        presentation = Presentation(f"{path}")
        text_runs = []

        for slide in presentation.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)

        file_text = " ".join(text_runs)

        if CUT_STR:
            file_text = file_text[:max_output_length]

        return file_text


        
if __name__ == "__main__":
    print("FileToText running...")
    # print(FileToText.pdf_to_text("../test.pdf", True))
    # print(FileToText.pdf_to_text("./test_files/1007_cw.pdf", CUT_STR=True, max_output_length=100))
    # print(FileToText.pdf_to_text("./test_files/1007_cw.pdf"))
    
    # print(FileToText.pptx_to_text("./test_files/COMP1004.pptx", CUT_STR=True, max_output_length=100))

    # print(FileToText.docx_to_text("./test_files/doc.docx"))
